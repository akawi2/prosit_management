from io import BytesIO
import os
import re

from django.http import HttpResponse, JsonResponse
from pptx import Presentation
from django.views.decorators.csrf import csrf_exempt
from prosit_management import settings

def remove_special_characters(input_string):
    # Regular expression to match non-alphanumeric characters (excluding spaces)
    cleaned_string = re.sub(r'[^a-zA-Z0-9\s]', '', input_string)
    return cleaned_string.strip()

@csrf_exempt
def downloadRetour(request):
    if request.method == 'POST':
        template = request.POST.get('template')
        if(template is None):
            print(template)
            return JsonResponse({'error': 'Please select a template.'}, status=400)
        
        titre = request.POST.get('titre')
        contexte = request.POST.get('contexte')
        date = request.POST.get('date')
        group = request.POST.get('group')

        problematique = request.POST.get('problematique')
        generalite = request.POST.get('generalite')
        besoins = []
        besoinStr = ""
        contraintes = []
        contraintesStr = ""
        pisteDeSolutions = []
        pisteDeSolutionsStr = ""
        planDActions = []
        planDActionStr = ""
        motCles = []
        definitions = []
        mots = ""
        motDefinitions = ""
        for key in request.POST:
            if key.startswith('besoin'):
                besoins.append(request.POST[key])
                besoinStr += f"- {request.POST[key]}\n"
                
        for key in request.POST:
            if key.startswith('mot'):
                motCles.append(request.POST[key])
                        
        for key in request.POST:
            if key.startswith('definition'):
                definitions.append(request.POST[key])
        
        for key in request.POST:
            if key.startswith('contrainte'):
                contraintes.append(request.POST[key])
                contraintesStr += f'- {request.POST[key]}\n' 
                    
        for key in request.POST:
            if key.startswith('piste'):
                pisteDeSolutions.append(request.POST[key])
                pisteDeSolutionsStr += f'- {request.POST[key]}\n'
        
        for key in request.POST:
            if key.startswith('plan'):
                planDActions.append(request.POST[key])
                planDActionStr += f'- {request.POST[key]}\n'
        
        
        for m,d in zip(motCles, definitions):
            mots += f"- {m}\n" 
            motDefinitions += f"- {m}: {d}\n"
            
        data = {
            '{{date}}': date,
            '{{group}}': group,
            '{{titre}}': titre,
            '{{mot_cles}}': mots,
            '{{mot_cles_definition}}': motDefinitions,
            '{{contexte}}': contexte,
            '{{problematique}}': problematique,
            '{{generalite}}': generalite,
            '{{besoin}}': besoinStr,
            '{{contrainte}}': contraintesStr,
            '{{piste_de_solution}}': pisteDeSolutionsStr,
            '{{plan_d_action}}': planDActionStr,
        }
        
        template_path = os.path.join(settings.BASE_DIR, 'prosit_management', 'doc_template', f'{template}.pptx')

        prs = Presentation(template_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for placeholder, value in data.items():
                        if placeholder in shape.text:
                            shape.text = shape.text.replace(placeholder, value)
    
        # Save the document to a response
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename="Prosit Retour {remove_special_characters(titre)}.pptx"'
        prs.save(response)  
        return response
